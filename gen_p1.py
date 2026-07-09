# -*- coding: utf-8 -*-
import os as _os; _os.chdir(_os.path.dirname(__file__))
from pptx import Presentation
from pptx.util import Inches as I, Pt
from pptx.dml.color import RGBColor as RC
from pptx.enum.text import PP_ALIGN as PA
from pptx.enum.shapes import MSO_SHAPE as MS

DB=RC(0x10,0x18,0x28); MB=RC(0x1D,0x4E,0xD8); LB=RC(0x3B,0x82,0xF6)
AB=RC(0x60,0xA5,0xFA); WH=RC(0xFF,0xFF,0xFF); LG=RC(0xF1,0xF5,0xF9)
DG=RC(0x33,0x40,0x55); MG=RC(0x64,0x74,0x8B); TD=RC(0x0F,0x17,0x2A)
GN=RC(0x10,0xB9,0x81); FN='Microsoft YaHei'; SW=I(13.333); SH=I(7.5)

p=Presentation(); p.slide_width,p.slide_height=SW,SH

def bg(s,c=WH):
    s.background.fill.solid(); s.background.fill.fore_color.rgb=c

def rc(s,l,t,w,h,fc=None,lc=None):
    r=s.shapes.add_shape(MS.RECTANGLE,l,t,w,h); r.line.fill.background()
    if fc: r.fill.solid(); r.fill.fore_color.rgb=fc
    else: r.fill.background()
    if lc: r.line.color.rgb=lc; r.line.width=Pt(1)

def tx(s,l,t,w,h,te,fs=14,c=TD,b=False,a=PA.LEFT):
    tb=s.shapes.add_textbox(l,t,w,h); tb.text_frame.word_wrap=True
    pp=tb.text_frame.paragraphs[0]; pp.text=te; pp.font.size=Pt(fs)
    pp.font.color.rgb=c; pp.font.bold=b; pp.font.name=FN; pp.alignment=a

def hdr(s,ti,st=''):
    rc(s,I(0),I(0),SW,I(1.3),fc=DB); rc(s,I(0),I(1.3),SW,I(0.06),fc=LB)
    tx(s,I(0.8),I(0.2),I(11),I(0.7),ti,fs=30,c=WH,b=True)
    if st: tx(s,I(0.8),I(0.82),I(11),I(0.45),st,fs=14,c=AB)

def ftr(s,pn):
    tx(s,I(11.5),I(7.0),I(1.5),I(0.4),str(pn),fs=11,c=MG,a=PA.RIGHT)
    rc(s,I(0),I(7.35),SW,I(0.15),fc=LB)

def bt(s,te,l=I(0.8),t=I(1.6)):
    tx(s,l,t,I(12),I(0.6),te,fs=22,c=DB,b=True)

def card(s,l,t,w,h,ti='',its=None,fc=WH,bc=LB):
    rc(s,l,t,w,h,fc=fc,lc=bc); yo=I(0.15)
    if ti: tx(s,l+I(0.2),t+yo,w-I(0.4),I(0.35),ti,fs=15,c=DB,b=True); yo+=I(0.4)
    if its:
        for it in its:
            tb=s.shapes.add_textbox(l+I(0.3),t+yo,w-I(0.6),I(0.28))
            pp=tb.text_frame.paragraphs[0]; pp.text=chr(8226)+' '+it
            pp.font.size=Pt(11); pp.font.color.rgb=DG; pp.font.name=FN
            yo+=I(0.26)

def tbl(s,l,t,cw,hd,rs,fs=10):
    nr,nc=len(rs)+1,len(hd); ts=s.shapes.add_table(nr,nc,l,t,sum(cw),I(0.35)*nr)
    tb=ts.table
    for ci,w in enumerate(cw): tb.columns[ci].width=w
    for ci,h in enumerate(hd):
        cl=tb.cell(0,ci); cl.text=h
        for pp in cl.text_frame.paragraphs:
            pp.font.size=Pt(fs); pp.font.bold=True; pp.font.color.rgb=WH
            pp.font.name=FN; pp.alignment=PA.CENTER
        cl.fill.solid(); cl.fill.fore_color.rgb=MB
    for ri,row in enumerate(rs):
        for ci,val in enumerate(row):
            cl=tb.cell(ri+1,ci); cl.text=str(val)
            for pp in cl.text_frame.paragraphs:
                pp.font.size=Pt(fs-1); pp.font.color.rgb=TD; pp.font.name=FN; pp.alignment=PA.CENTER
            cl.fill.solid(); cl.fill.fore_color.rgb=LG if ri%2==0 else WH

def arr(s,l,t,w=I(0.45)):
    a=s.shapes.add_shape(MS.RIGHT_ARROW,l,t,w,I(0.22))
    a.fill.solid(); a.fill.fore_color.rgb=LB; a.line.fill.background()
